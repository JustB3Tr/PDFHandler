import io
import zipfile
from dataclasses import dataclass

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


@dataclass
class RenderedPage:
    page_number: int
    image: Image.Image
    width_in: float
    height_in: float


def render_pdf_preview(pdf_path: str, preview_dpi: int = 110):
    """Render lightweight source previews."""
    doc = fitz.open(pdf_path)
    pages = []

    for i in range(doc.page_count):
        page = doc[i]
        rect = page.rect
        pix = page.get_pixmap(dpi=preview_dpi, alpha=False)
        img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")

        pages.append(
            RenderedPage(
                page_number=i + 1,
                image=img,
                width_in=rect.width / 72,
                height_in=rect.height / 72,
            )
        )

    return pages


def render_pdf_full(pdf_path: str, dpi: int, progress_callback=None):
    """Render full-quality pages for export."""
    doc = fitz.open(pdf_path)
    pages = []

    for i in range(doc.page_count):
        page = doc[i]
        rect = page.rect
        pix = page.get_pixmap(dpi=dpi, alpha=False)
        img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")

        rendered = RenderedPage(
            page_number=i + 1,
            image=img,
            width_in=rect.width / 72,
            height_in=rect.height / 72,
        )
        pages.append(rendered)

        if progress_callback:
            progress_callback(i + 1, doc.page_count, rendered)

    return pages


def _new_clean_presentation(width_in: float, height_in: float):
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)

    # Remove default starter slide cleanly
    if len(prs.slides) > 0:
        slide_ids = prs.slides._sldIdLst
        while len(slide_ids) > 0:
            rId = slide_ids[0].rId
            prs.part.drop_rel(rId)
            del slide_ids[0]

    return prs


def export_pptx(rendered_pages, out_path: str):
    if not rendered_pages:
        raise ValueError("No pages to export.")

    first = rendered_pages[0]
    prs = _new_clean_presentation(first.width_in, first.height_in)
    blank_layout = prs.slide_layouts[6]

    for page in rendered_pages:
        slide = prs.slides.add_slide(blank_layout)
        stream = io.BytesIO()
        page.image.save(stream, format="PNG")
        stream.seek(0)

        slide.shapes.add_picture(
            stream,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(out_path)


def export_png_zip(rendered_pages, out_path: str):
    if not rendered_pages:
        raise ValueError("No pages to export.")

    with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for page in rendered_pages:
            stream = io.BytesIO()
            page.image.save(stream, format="PNG")
            zf.writestr(f"slide_{page.page_number:03d}.png", stream.getvalue())


def export_jpg_zip(rendered_pages, out_path: str):
    if not rendered_pages:
        raise ValueError("No pages to export.")

    with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for page in rendered_pages:
            stream = io.BytesIO()
            page.image.convert("RGB").save(stream, format="JPEG", quality=95)
            zf.writestr(f"slide_{page.page_number:03d}.jpg", stream.getvalue())


EXPORT_HANDLERS = {
    "pptx": export_pptx,
    "gslides_pptx": export_pptx,  # same export, different label in UI
    "png_zip": export_png_zip,
    "jpg_zip": export_jpg_zip,
}


def default_extension_for(fmt: str) -> str:
    return {
        "pptx": ".pptx",
        "gslides_pptx": ".pptx",
        "png_zip": ".zip",
        "jpg_zip": ".zip",
    }[fmt]


def suggested_filetypes(fmt: str):
    return {
        "pptx": [("PowerPoint Presentation", "*.pptx")],
        "gslides_pptx": [("PowerPoint Presentation", "*.pptx")],
        "png_zip": [("ZIP Archive", "*.zip")],
        "jpg_zip": [("ZIP Archive", "*.zip")],
    }[fmt]